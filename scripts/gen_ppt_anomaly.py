"""
多模态生成精度异常判定 - PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色方案 ──
C_BG = RGBColor(0xF5, 0xF7, 0xFA)          # 页面底色
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)        # 深色文字
C_PRIMARY = RGBColor(0x0D, 0x47, 0xA1)     # 主蓝
C_PRIMARY_L = RGBColor(0x1E, 0x88, 0xE5)   # 亮蓝
C_PRIMARY_XL = RGBColor(0xE3, 0xF2, 0xFD)  # 浅蓝底
C_ACCENT_O = RGBColor(0xE6, 0x5C, 0x00)    # 橙色强调
C_ACCENT_G = RGBColor(0x2E, 0x7D, 0x32)    # 绿色
C_ACCENT_R = RGBColor(0xC6, 0x28, 0x28)    # 红色
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x75, 0x75, 0x75)
C_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
C_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_BANNER_BG = RGBColor(0x0D, 0x47, 0xA1)   # 横幅深蓝
C_TARGET_BG = RGBColor(0xE6, 0x5C, 0x00)   # 目标横幅橙
C_BOTTLENECK1 = RGBColor(0xE8, 0xEA, 0xF6) # 瓶颈卡片1 淡靛蓝
C_BOTTLENECK2 = RGBColor(0xFE, 0xF3, 0xE2) # 瓶颈卡片2 淡橙
C_BOTTLENECK3 = RGBColor(0xFC, 0xE4, 0xEC) # 瓶颈卡片3 淡红
C_ARCH_NODE = RGBColor(0x0D, 0x47, 0xA1)   # 架构节点深蓝
C_ARCH_NODE_L = RGBColor(0x1E, 0x88, 0xE5) # 架构节点亮蓝
C_ARCH_L1 = RGBColor(0x1B, 0x5E, 0x20)     # L1绿
C_ARCH_L2 = RGBColor(0xE6, 0x5C, 0x00)     # L2橙
C_ARCH_L3 = RGBColor(0x6A, 0x1B, 0x9A)     # L3紫
C_ADV_BG1 = RGBColor(0xE8, 0xF5, 0xE9)     # 优势卡片绿底
C_ADV_BG2 = RGBColor(0xE3, 0xF2, 0xFD)     # 优势卡片蓝底
C_ADV_BG3 = RGBColor(0xF3, 0xE5, 0xF5)     # 优势卡片紫底


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    # 圆角半径
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=12, bold=False, color=C_DARK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """设置形状文本"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_text_para(tf, text, font_size=12, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
                  space_before=0, space_after=0, font_name="微软雅黑"):
    """在已有 text_frame 上追加段落"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                color=C_DARK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_arrow_down(slide, cx, top, length, color=C_PRIMARY_L):
    """添加向下箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, cx - Inches(0.12), top, Inches(0.24), length
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build_slide(prs):
    """构建唯一的一页 PPT"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # ── 页面背景 ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG

    SW = prs.slide_width   # 13.333 inches
    SH = prs.slide_height  # 7.5 inches
    MARGIN = Inches(0.3)

    # ============================================================
    # 顶部区域：标题 + 背景横幅 + 目标横幅
    # ============================================================

    # 标题
    add_textbox(slide, MARGIN, Inches(0.15), Inches(8), Inches(0.45),
                "多模态生成精度异常判定", font_size=22, bold=True, color=C_PRIMARY)

    # 背景横幅
    banner_h = Inches(0.55)
    banner = add_rounded_rect(slide, MARGIN, Inches(0.55), SW - 2 * MARGIN, banner_h,
                              C_BANNER_BG)
    tf = set_text(banner, "", font_size=10)
    tf.paragraphs[0].clear()
    # 使用 run 来做加粗标签 + 正文
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = " 背景  "
    run1.font.size = Pt(10)
    run1.font.bold = True
    run1.font.color.rgb = C_WHITE
    run1.font.name = "微软雅黑"
    run2 = p.add_run()
    run2.text = "受限于生成模型的精度瓶颈，AIGC 内容常伴随手指畸形、运动反向等细粒度生成误差，现有评估缺乏对精度缺陷的量化能力，亟需高精度针对性检测方案"
    run2.font.size = Pt(10)
    run2.font.color.rgb = C_WHITE
    run2.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.LEFT
    banner.text_frame.margin_left = Inches(0.15)
    banner.text_frame.margin_top = Inches(0.05)

    # 目标横幅
    target = add_rounded_rect(slide, MARGIN, Inches(1.15), SW - 2 * MARGIN, Inches(0.38),
                              C_TARGET_BG)
    tf2 = set_text(target, "", font_size=10)
    p2 = tf2.paragraphs[0]
    r1 = p2.add_run()
    r1.text = " 目标  "
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = C_WHITE
    r1.font.name = "微软雅黑"
    r2 = p2.add_run()
    r2.text = "构建多模态异常检测系统，各类业务典型异常识别准召率达成 90+%"
    r2.font.size = Pt(10)
    r2.font.color.rgb = C_WHITE
    r2.font.name = "微软雅黑"
    target.text_frame.margin_left = Inches(0.15)
    target.text_frame.margin_top = Inches(0.02)

    # ============================================================
    # 分栏起始 Y
    # ============================================================
    COL_TOP = Inches(1.65)
    LEFT_W = Inches(6.1)
    RIGHT_X = Inches(6.7)
    RIGHT_W = Inches(6.3)

    # ============================================================
    # 左半部分
    # ============================================================

    # ── 小标题：问题与挑战 ──
    left_title = add_textbox(slide, MARGIN, COL_TOP, LEFT_W, Inches(0.3),
                             "▎问题与挑战", font_size=13, bold=True, color=C_PRIMARY)

    # ── 异常帧定位 ──
    loc_y = COL_TOP + Inches(0.35)
    loc_card = add_rounded_rect(slide, MARGIN, loc_y, LEFT_W, Inches(0.45),
                                C_PRIMARY_XL, border_color=C_PRIMARY_L, border_width=1)
    tf_loc = set_text(loc_card, "", font_size=9)
    p_loc = tf_loc.paragraphs[0]
    r_icon = p_loc.add_run()
    r_icon.text = ">> 异常帧定位  "
    r_icon.font.size = Pt(9.5)
    r_icon.font.bold = True
    r_icon.font.color.rgb = C_PRIMARY
    r_icon.font.name = "微软雅黑"
    r_desc = p_loc.add_run()
    r_desc.text = "对检出的异常图像/视频，精准标记异常起始帧，实现问题的快速溯源与定位"
    r_desc.font.size = Pt(9)
    r_desc.font.color.rgb = C_DARK
    r_desc.font.name = "微软雅黑"
    loc_card.text_frame.margin_left = Inches(0.12)
    loc_card.text_frame.margin_top = Inches(0.05)

    # ── 异常图谱表格占位 ──
    table_y = loc_y + Inches(0.55)
    add_textbox(slide, MARGIN + Inches(0.05), table_y, LEFT_W, Inches(0.22),
                "异常图谱多元复杂，精准识别与系统性治理难度高：", font_size=9, bold=True, color=C_DARK)

    # 典型异常表格
    tbl_y = table_y + Inches(0.22)
    tbl_w = LEFT_W - Inches(0.1)
    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(rows, cols, MARGIN + Inches(0.05), tbl_y, tbl_w, Inches(1.2))
    table = table_shape.table

    # 列宽
    table.columns[0].width = Inches(1.1)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(1.4)

    # 表头
    headers = ["异常大类", "典型异常", "表现描述", "检测难度"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY

    # 数据行
    data = [
        ["生物结构", "手指畸形/多余", "手指数量异常、关节角度超限、骨骼比例突变", "中"],
        ["生物结构", "面部/眼部异常", "非自然眨眼频率、瞳孔追踪失败、口腔纹理伪影", "高"],
        ["运动逻辑", "运动反向/突变", "光流加速度跳变、物体运动违反惯性", "高"],
        ["物理/背景", "物理常识违反", "像素漂移、重力方向不一致、背景深度闪烁", "极高"],
    ]
    row_colors = [C_WHITE, RGBColor(0xF5, 0xF5, 0xF5)]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7.5)
                p.font.color.rgb = C_DARK
                p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.CENTER if c in [0, 3] else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_colors[r % 2]
            # 难度列着色
            if c == 3:
                if val == "极高":
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = C_ACCENT_R
                        p.font.bold = True
                elif val == "高":
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = C_ACCENT_O
                        p.font.bold = True

    # ── 三大瓶颈 ──
    bn_label_y = tbl_y + Inches(1.30)
    add_textbox(slide, MARGIN + Inches(0.05), bn_label_y, LEFT_W, Inches(0.22),
                "评估体系三大瓶颈：", font_size=9, bold=True, color=C_ACCENT_R)

    bn_y = bn_label_y + Inches(0.25)
    bn_w = Inches(1.93)
    bn_h = Inches(2.35)
    bn_gap = Inches(0.1)

    bottleneck_data = [
        {
            "title": "指标失准",
            "bg": C_BOTTLENECK1,
            "icon_color": C_PRIMARY,
            "items": [
                "客观指标与主观感知不一致，FVD、IS-V 等基于特定特征空间计算，与人类视觉判断产生系统性偏差",
                "单一指标无法全面反映综合质量",
            ]
        },
        {
            "title": "真值稀缺",
            "bg": C_BOTTLENECK2,
            "icon_color": C_ACCENT_O,
            "items": [
                "生成内容无 Ground Truth，现有无参考指标对生成伪影的 heavy-tailed 分布拟合不足",
                "主观标注成本高昂，专业评估者与普适用户间存在感知差异",
            ]
        },
        {
            "title": "标准缺失",
            "bg": C_BOTTLENECK3,
            "icon_color": C_ACCENT_R,
            "items": [
                "行业统一标准尚未建立",
                "生成模型发展过快，检测技术发展相对滞后",
            ]
        },
    ]

    for i, bn in enumerate(bottleneck_data):
        bx = MARGIN + i * (bn_w + bn_gap)
        card = add_rounded_rect(slide, bx, bn_y, bn_w, bn_h, bn["bg"])
        card.text_frame.margin_left = Inches(0.1)
        card.text_frame.margin_right = Inches(0.08)
        card.text_frame.margin_top = Inches(0.08)

        tf_bn = card.text_frame
        tf_bn.word_wrap = True
        # 标题
        p_title = tf_bn.paragraphs[0]
        p_title.text = f"  {bn['title']}"
        p_title.font.size = Pt(10)
        p_title.font.bold = True
        p_title.font.color.rgb = bn["icon_color"]
        p_title.font.name = "微软雅黑"
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(6)

        # 分隔线（用文字模拟）
        p_sep = tf_bn.add_paragraph()
        p_sep.text = "─" * 18
        p_sep.font.size = Pt(5)
        p_sep.font.color.rgb = C_LIGHT_GRAY
        p_sep.alignment = PP_ALIGN.CENTER
        p_sep.space_after = Pt(4)

        for item_text in bn["items"]:
            p_item = tf_bn.add_paragraph()
            p_item.text = f"• {item_text}"
            p_item.font.size = Pt(7)
            p_item.font.color.rgb = C_DARK
            p_item.font.name = "微软雅黑"
            p_item.space_before = Pt(3)
            p_item.space_after = Pt(2)

    # ============================================================
    # 右半部分
    # ============================================================

    # ── 小标题：检测系统方案 ──
    add_textbox(slide, RIGHT_X, COL_TOP, RIGHT_W, Inches(0.3),
                "▎检测系统方案（图像 + 视频双模态）", font_size=13, bold=True, color=C_PRIMARY)

    # ── 架构流程图 ──
    arch_y = COL_TOP + Inches(0.35)
    node_h = Inches(0.35)
    node_w = Inches(5.8)
    node_x = RIGHT_X + Inches(0.25)
    arrow_len = Inches(0.22)
    gap = Inches(0.04)

    # --- 输入层 ---
    n1 = add_rounded_rect(slide, node_x, arch_y, node_w, node_h, C_ARCH_NODE)
    tf1 = set_text(n1, "输入层：图像 / 视频帧序列", font_size=9, bold=True, color=C_WHITE,
                   align=PP_ALIGN.CENTER)
    n1.text_frame.margin_top = Inches(0.02)

    # 箭头1
    a1_y = arch_y + node_h + gap
    add_arrow_down(slide, node_x + node_w // 2, a1_y, arrow_len)

    # --- 特征提取层 ---
    feat_y = a1_y + arrow_len + gap
    feat_h = Inches(0.7)
    feat_card = add_rounded_rect(slide, node_x, feat_y, node_w, feat_h,
                                  C_PRIMARY_XL, border_color=C_PRIMARY_L, border_width=1)
    tf_feat = feat_card.text_frame
    tf_feat.word_wrap = True
    tf_feat.margin_left = Inches(0.1)
    tf_feat.margin_top = Inches(0.03)
    p_ft = tf_feat.paragraphs[0]
    p_ft.text = "共享特征提取层（Feature Hub）"
    p_ft.font.size = Pt(9)
    p_ft.font.bold = True
    p_ft.font.color.rgb = C_PRIMARY
    p_ft.font.name = "微软雅黑"
    p_ft.alignment = PP_ALIGN.CENTER

    feat_items = [
        "关键点检测（MediaPipe Face/Hand/Body）    人脸嵌入（InsightFace ArcFace）",
        "深度估计（MiDaS）    光流分析（RAFT，仅视频）    纹理/频域特征",
    ]
    for fi in feat_items:
        pf = tf_feat.add_paragraph()
        pf.text = fi
        pf.font.size = Pt(7)
        pf.font.color.rgb = C_GRAY
        pf.font.name = "微软雅黑"
        pf.alignment = PP_ALIGN.CENTER
        pf.space_before = Pt(1)

    # 箭头2
    a2_y = feat_y + feat_h + gap
    add_arrow_down(slide, node_x + node_w // 2, a2_y, arrow_len)

    # --- 三级检测引擎 ---
    eng_y = a2_y + arrow_len + gap
    eng_h = Inches(1.05)
    eng_card = add_rounded_rect(slide, node_x, eng_y, node_w, eng_h,
                                 C_CARD_BG, border_color=C_PRIMARY_L, border_width=1.5)
    tf_eng = eng_card.text_frame
    tf_eng.word_wrap = True
    tf_eng.margin_left = Inches(0.1)
    tf_eng.margin_top = Inches(0.03)
    p_eng_title = tf_eng.paragraphs[0]
    p_eng_title.text = "多级异常检测引擎（三级递进）"
    p_eng_title.font.size = Pt(9)
    p_eng_title.font.bold = True
    p_eng_title.font.color.rgb = C_PRIMARY
    p_eng_title.font.name = "微软雅黑"
    p_eng_title.alignment = PP_ALIGN.CENTER

    levels = [
        ("L1 快速规则筛查", "几何约束（手指数量/关节角度/骨骼比例）、生理指标（EAR/MAR）", C_ARCH_L1),
        ("L2 结构化分析", "骨架拓扑验证、纹理一致性、时序平滑度（图像+视频）", C_ARCH_L2),
        ("L3 MLLM 语义推理", "大模型视觉理解兜底，覆盖长尾异常与复合型缺陷", C_ARCH_L3),
    ]
    for lvl_name, lvl_desc, lvl_color in levels:
        p_lvl = tf_eng.add_paragraph()
        r_name = p_lvl.add_run()
        r_name.text = f"  {lvl_name}  "
        r_name.font.size = Pt(8)
        r_name.font.bold = True
        r_name.font.color.rgb = lvl_color
        r_name.font.name = "微软雅黑"
        r_d = p_lvl.add_run()
        r_d.text = lvl_desc
        r_d.font.size = Pt(7)
        r_d.font.color.rgb = C_DARK
        r_d.font.name = "微软雅黑"
        p_lvl.space_before = Pt(3)

    # 箭头3
    a3_y = eng_y + eng_h + gap
    add_arrow_down(slide, node_x + node_w // 2, a3_y, arrow_len)

    # --- 六维度评估矩阵 ---
    dim_y = a3_y + arrow_len + gap
    dim_h = Inches(1.1)
    dim_card = add_rounded_rect(slide, node_x, dim_y, node_w, dim_h,
                                 C_PRIMARY_XL, border_color=C_PRIMARY_L, border_width=1)
    tf_dim = dim_card.text_frame
    tf_dim.word_wrap = True
    tf_dim.margin_left = Inches(0.08)
    tf_dim.margin_top = Inches(0.03)
    p_dim_title = tf_dim.paragraphs[0]
    p_dim_title.text = "六维度评估矩阵"
    p_dim_title.font.size = Pt(9)
    p_dim_title.font.bold = True
    p_dim_title.font.color.rgb = C_PRIMARY
    p_dim_title.font.name = "微软雅黑"
    p_dim_title.alignment = PP_ALIGN.CENTER

    dimensions = [
        ("D1 人脸一致性", "CSIM 余弦相似度", "D2 表情自然度", "AU + FACS 规则"),
        ("D3 生物结构异常", "眼/手/口/躯体", "D4 运动逻辑", "光流加速度平滑性"),
        ("D5 物理常识", "像素漂移+重力一致性", "D6 背景一致性", "深度时序+特征匹配"),
    ]
    for d1_name, d1_desc, d2_name, d2_desc in dimensions:
        p_row = tf_dim.add_paragraph()
        r1n = p_row.add_run()
        r1n.text = f"  {d1_name}"
        r1n.font.size = Pt(7.5)
        r1n.font.bold = True
        r1n.font.color.rgb = C_PRIMARY
        r1n.font.name = "微软雅黑"
        r1d = p_row.add_run()
        r1d.text = f" ({d1_desc})    "
        r1d.font.size = Pt(7)
        r1d.font.color.rgb = C_GRAY
        r1d.font.name = "微软雅黑"
        r2n = p_row.add_run()
        r2n.text = f"  {d2_name}"
        r2n.font.size = Pt(7.5)
        r2n.font.bold = True
        r2n.font.color.rgb = C_PRIMARY
        r2n.font.name = "微软雅黑"
        r2d = p_row.add_run()
        r2d.text = f" ({d2_desc})"
        r2d.font.size = Pt(7)
        r2d.font.color.rgb = C_GRAY
        r2d.font.name = "微软雅黑"
        p_row.space_before = Pt(2)

    # 箭头4
    a4_y = dim_y + dim_h + gap
    add_arrow_down(slide, node_x + node_w // 2, a4_y, arrow_len)

    # --- 输出层 ---
    out_y = a4_y + arrow_len + gap
    out_h = Inches(0.48)
    out_card = add_rounded_rect(slide, node_x, out_y, node_w, out_h, C_ARCH_NODE)
    tf_out = out_card.text_frame
    tf_out.word_wrap = True
    tf_out.margin_left = Inches(0.1)
    tf_out.margin_top = Inches(0.02)
    p_out = tf_out.paragraphs[0]
    p_out.text = "输出层"
    p_out.font.size = Pt(9)
    p_out.font.bold = True
    p_out.font.color.rgb = C_WHITE
    p_out.font.name = "微软雅黑"
    p_out.alignment = PP_ALIGN.CENTER
    p_out2 = tf_out.add_paragraph()
    p_out2.text = "异常类型分类 + 置信度评分  |  异常帧精准定位（起始帧标记）  |  多维度加权综合评分"
    p_out2.font.size = Pt(7)
    p_out2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p_out2.font.name = "微软雅黑"
    p_out2.alignment = PP_ALIGN.CENTER

    # ── 方案核心优势（三个小卡片） ──
    adv_y = out_y + out_h + Inches(0.12)
    adv_w = Inches(1.87)
    adv_h = Inches(0.7)
    adv_gap = Inches(0.1)

    advantages = [
        {
            "title": "双模态覆盖",
            "desc": "统一特征提取，兼顾静态图像异常与动态视频时序异常",
            "bg": C_ADV_BG1,
            "color": C_ACCENT_G,
        },
        {
            "title": "三级递进检测",
            "desc": "规则→结构→语义，兼顾效率与覆盖率，逐层收敛",
            "bg": C_ADV_BG2,
            "color": C_PRIMARY,
        },
        {
            "title": "MLLM 兜底长尾",
            "desc": "大模型语义理解覆盖规则难以枚举的异常类型",
            "bg": C_ADV_BG3,
            "color": C_ARCH_L3,
        },
    ]

    for i, adv in enumerate(advantages):
        ax = node_x + i * (adv_w + adv_gap)
        acard = add_rounded_rect(slide, ax, adv_y, adv_w, adv_h, adv["bg"])
        tf_a = acard.text_frame
        tf_a.word_wrap = True
        tf_a.margin_left = Inches(0.08)
        tf_a.margin_top = Inches(0.05)
        pa_t = tf_a.paragraphs[0]
        pa_t.text = adv["title"]
        pa_t.font.size = Pt(8.5)
        pa_t.font.bold = True
        pa_t.font.color.rgb = adv["color"]
        pa_t.font.name = "微软雅黑"
        pa_t.alignment = PP_ALIGN.CENTER
        pa_d = tf_a.add_paragraph()
        pa_d.text = adv["desc"]
        pa_d.font.size = Pt(6.5)
        pa_d.font.color.rgb = C_GRAY
        pa_d.font.name = "微软雅黑"
        pa_d.alignment = PP_ALIGN.CENTER
        pa_d.space_before = Pt(3)

    # ── 中间分隔线 ──
    sep_x = Inches(6.45)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sep_x, COL_TOP, Inches(0.02), Inches(5.6))
    sep.fill.solid()
    sep.fill.fore_color.rgb = C_LIGHT_GRAY
    sep.line.fill.background()

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide(prs)

    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "outputs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "多模态生成精度异常判定.pptx")
    prs.save(out_path)
    print(f"PPT 已生成: {out_path}")


if __name__ == "__main__":
    main()
